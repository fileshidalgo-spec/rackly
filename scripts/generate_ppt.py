#!/usr/bin/env python3
"""
RACKLY — Manual de Usuario Profesional
Presentación minimalista · Formato APA · Gramática RAE
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Paleta de colores minimalista ───
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x2E)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MED_GRAY    = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY  = RGBColor(0xF0, 0xF0, 0xF5)
ACCENT      = RGBColor(0x3B, 0x82, 0xF6)   # Azul corporativo
ACCENT_DARK = RGBColor(0x1E, 0x40, 0xAF)
ACCENT_LT   = RGBColor(0xDB, 0xEA, 0xFE)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
RED_ACCENT  = RGBColor(0xEF, 0x44, 0x44)
ORANGE      = RGBColor(0xF5, 0x9E, 0x0B)
TEAL        = RGBColor(0x14, 0xB8, 0xA6)
PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)
ROSE        = RGBColor(0xF4, 0x3F, 0x5E)

# ─── Fuentes ───
FONT_MAIN = 'Carlito'
FONT_TITLE = 'Carlito'

# ─── Configuración de diapositiva ───
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ─── Helpers ───

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_GRAY, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_MAIN, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=14,
                          color=DARK_GRAY, bold=False, font_name=FONT_MAIN,
                          line_spacing=1.4, alignment=PP_ALIGN.LEFT):
    """lines: list of (text, bold, color) tuples or just strings"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, b, c = line_data, bold, color
        else:
            text = line_data[0]
            b = line_data[1] if len(line_data) > 1 else bold
            c = line_data[2] if len(line_data) > 2 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_circle_icon(slide, left, top, size, color, text='', font_size=14, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = FONT_MAIN
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape

def add_numbered_step(slide, left, top, num, title, description, accent_color=ACCENT,
                     title_size=16, desc_size=13):
    """Add a numbered step with circle icon + title + description."""
    circle_size = Inches(0.45)
    add_circle_icon(slide, left, top + Pt(2), circle_size, accent_color, str(num), 16, WHITE)
    add_textbox(slide, left + circle_size + Inches(0.2), top, Inches(5), Inches(0.35),
                title, title_size, BLACK, bold=True)
    add_textbox(slide, left + circle_size + Inches(0.2), top + Inches(0.32),
                Inches(5.2), Inches(0.8), description, desc_size, MED_GRAY)

def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=DARK_GRAY, bullet_char='\u2022', line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(item, tuple):
            text, b, c = item[0], item[1] if len(item) > 1 else False, item[2] if len(item) > 2 else color
        else:
            text, b, c = item, False, color
        p.text = f'{bullet_char}  {text}'
        p.font.size = Pt(font_size)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = FONT_MAIN
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_footer(slide, slide_num, total_slides):
    """Minimalista footer con número de diapositiva."""
    add_textbox(slide, Inches(0.5), SLIDE_H - Inches(0.5), Inches(6), Inches(0.3),
                'RACKLY v2.0 \u2014 Manual de Usuario', 9, MED_GRAY, italic=True)
    add_textbox(slide, SLIDE_W - Inches(2), SLIDE_H - Inches(0.5), Inches(1.5), Inches(0.3),
                f'{slide_num} / {total_slides}', 9, MED_GRAY, alignment=PP_ALIGN.RIGHT)

def make_section_slide(title, subtitle, accent_color=ACCENT, section_num=''):
    """Diapositiva de sección con fondo de color."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    # Barra lateral
    add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, accent_color)
    # Número de sección
    if section_num:
        add_textbox(slide, Inches(1.5), Inches(1.8), Inches(2), Inches(0.6),
                    section_num, 52, accent_color, bold=True)
    # Título
    add_textbox(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(0.8),
                title, 36, BLACK, bold=True)
    add_accent_line(slide, Inches(1.5), Inches(3.45), Inches(1.5), accent_color)
    # Subtítulo
    add_textbox(slide, Inches(1.5), Inches(3.7), Inches(9), Inches(0.5),
                subtitle, 16, MED_GRAY, italic=True)
    return slide

TOTAL_SLIDES = 20

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 1 — PORTADA
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
# Barra superior de acento
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
# Logo/branding area
add_circle_icon(slide, Inches(1.2), Inches(2.0), Inches(1.0), ACCENT, 'R', 36, WHITE)
add_textbox(slide, Inches(2.5), Inches(2.05), Inches(4), Inches(0.6),
            'RACKLY', 44, BLACK, bold=True)
add_textbox(slide, Inches(2.5), Inches(2.6), Inches(6), Inches(0.4),
            'Sistema de Gesti\u00f3n de Almacenes v2.0', 18, MED_GRAY)
# Línea divisoria
add_accent_line(slide, Inches(1.2), Inches(3.4), Inches(8), LIGHT_GRAY)
# Subtítulo principal
add_textbox(slide, Inches(1.2), Inches(3.8), Inches(10), Inches(0.7),
            'Manual de Usuario: Gu\u00eda Paso a Paso', 28, BLACK, bold=True)
add_textbox(slide, Inches(1.2), Inches(4.5), Inches(10), Inches(0.5),
            'Funcionalidades, flujos de trabajo y mejores pr\u00e1cticas', 16, MED_GRAY, italic=True)
# Datos de autor
add_multiline_textbox(slide, Inches(1.2), Inches(5.8), Inches(6), Inches(1.2), [
    ('Autor: Miguel Hidalgo', False, MED_GRAY),
    ('Aplicaci\u00f3n web \u2014 Next.js + Supabase', False, MED_GRAY),
    ('Junio de 2026', False, MED_GRAY),
], 13, line_spacing=1.6)
# Barra inferior
add_shape_bg(slide, Inches(0), SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), ACCENT)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 2 — CONTENIDO
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
            'Contenido', 32, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.2), ACCENT)

sections = [
    ('01', 'Introducci\u00f3n a RACKLY', 'Visi\u00f3n general y arquitectura del sistema'),
    ('02', 'Autenticaci\u00f3n y Roles', 'Registro, inicio de sesi\u00f3n y permisos'),
    ('03', 'Kardex Racks \u2014 Movimientos', 'Ingreso, salida, devoluci\u00f3n e INC'),
    ('04', 'Kardex Racks \u2014 Traslado', 'Transferencia entre ubicaciones'),
    ('05', 'Kardex Racks \u2014 Cat\u00e1logo', 'Gesti\u00f3n del cat\u00e1logo de productos'),
    ('06', 'Kardex Racks \u2014 Stock', 'Consulta de stock por c\u00f3digo'),
    ('07', 'Kardex Racks \u2014 Ocupaci\u00f3n', 'Mapa visual de posiciones'),
    ('08', 'Kardex Racks \u2014 FEFO', 'Control de fechas de vencimiento'),
    ('09', 'Kardex Racks \u2014 Descarga', 'Exportaci\u00f3n a Excel'),
    ('10', 'Kardex Racks \u2014 Usuarios', 'Administraci\u00f3n de personal'),
    ('11', 'Kardex Piso', 'Operaciones en piso de almac\u00e9n'),
    ('12', 'Funciones Transversales', 'Modo offline, sincronizaci\u00f3n y tiempo real'),
    ('13', 'Referencias', 'Fuentes en formato APA'),
]
y = Inches(1.3)
for num, title, desc in sections:
    add_textbox(slide, Inches(1.2), y, Inches(0.7), Inches(0.3), num, 14, ACCENT, bold=True)
    add_textbox(slide, Inches(2.0), y, Inches(4), Inches(0.3), title, 15, BLACK, bold=True)
    add_textbox(slide, Inches(2.0), y + Inches(0.25), Inches(6), Inches(0.25), desc, 11, MED_GRAY, italic=True)
    y += Inches(0.46)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 3 — SECCIÓN 01: Introducción
# ══════════════════════════════════════════════════════════════
make_section_slide('Introducci\u00f3n a RACKLY',
                   'Descripci\u00f3n general, prop\u00f3sito y arquitectura del sistema.',
                   ACCENT, '01')

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Introducci\u00f3n a RACKLY', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), ACCENT)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.9),
            'RACKLY es un sistema de gesti\u00f3n de almacenes dise\u00f1ado para controlar de forma precisa '
            'los movimientos de insumos, el stock disponible, las fechas de vencimiento y la ocupaci\u00f3n '
            'f\u00edsica de las posiciones de racks y piso. La aplicaci\u00f3n opera en tiempo real con capacidad '
            'offline, lo que garantiza continuidad operativa incluso sin conexi\u00f3n a internet.',
            14, DARK_GRAY, line_spacing=1.5)

# Características clave en tarjetas
cards = [
    ('Gesti\u00f3n de Movimientos', 'Registro de ingresos, salidas, devoluciones, traslados e insumos no conformes (INC).', ACCENT),
    ('Control de Stock', 'C\u00e1lculo autom\u00e1tico del stock por c\u00f3digo y ubicaci\u00f3n, con integraci\u00f3n al sistema Big Magic.', GREEN),
    ('FEFO', 'Primer vencido, primera salida: control de fechas con indicadores por colores.', ORANGE),
    ('Modo Offline', 'Funcionamiento sin internet con sincronizaci\u00f3n autom\u00e1tica al reconectar.', PURPLE),
    ('Roles y Permisos', 'Administraci\u00f3n jer\u00e1rquica de usuarios: operarios, supervisores y administradores.', TEAL),
    ('Mapa Visual', 'Vista tridimensional de la ocupaci\u00f3n del almac\u00e9n con actualizaci\u00f3n en tiempo real.', ROSE),
]
x = Inches(0.8)
y = Inches(2.7)
for i, (title, desc, color) in enumerate(cards):
    col = i % 3
    row = i // 3
    cx = x + col * Inches(4.1)
    cy = y + row * Inches(2.1)
    card_bg = add_shape_bg(slide, cx, cy, Inches(3.8), Inches(1.8), LIGHT_GRAY)
    add_shape_bg(slide, cx, cy, Inches(0.06), Inches(1.8), color)
    add_textbox(slide, cx + Inches(0.2), cy + Inches(0.15), Inches(3.4), Inches(0.3),
                title, 14, BLACK, bold=True)
    add_textbox(slide, cx + Inches(0.2), cy + Inches(0.5), Inches(3.4), Inches(1.1),
                desc, 12, MED_GRAY, line_spacing=1.4)

add_footer(slide, 3, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 4 — SECCIÓN 02: Autenticación y Roles
# ══════════════════════════════════════════════════════════════
make_section_slide('Autenticaci\u00f3n y Roles de Usuario',
                   'Registro, inicio de sesi\u00f3n, aprobaci\u00f3n y jerarqu\u00eda de permisos.',
                   TEAL, '02')

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, TEAL)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Flujo de Autenticaci\u00f3n', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), TEAL)

steps_auth = [
    ('Registro de cuenta', 'El usuario solicita acceso con correo electr\u00f3nico y contrase\u00f1a (m\u00ednimo 6 caracteres). El sistema confirma autom\u00e1ticamente el correo sin necesidad de verificaci\u00f3n manual.'),
    ('Aprobaci\u00f3n administrativa', 'Las cuentas nuevas permanecen en estado \u00abpendiente\u00bb hasta que un administrador las apruebe desde la pesta\u00f1a \u00abUsuarios\u00bb. Sin aprobaci\u00f3n, el acceso al sistema queda restringido.'),
    ('Inicio de sesi\u00f3n', 'Tras la aprobaci\u00f3n, el usuario ingresa con sus credenciales. Si es su primer inicio, puede ser requerido cambiar la contrase\u00f1a seg\u00fan pol\u00edticas de seguridad.'),
    ('Asignaci\u00f3n de rol', 'El administrador asigna un rol espec\u00edfico que define los permisos del usuario dentro del sistema. Los roles se pueden modificar en cualquier momento.'),
]
y = Inches(1.4)
for i, (title, desc) in enumerate(steps_auth):
    add_numbered_step(slide, Inches(0.8), y, i+1, title, desc, TEAL)
    y += Inches(1.25)

add_footer(slide, 4, TOTAL_SLIDES)

# Diapositiva 5 — Roles
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, TEAL)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Jerarqu\u00eda de Roles', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), TEAL)

roles = [
    ('Administrador', 'Acceso completo: gesti\u00f3n de usuarios, movimientos, cat\u00e1logo, configuraci\u00f3n y eliminaci\u00f3n de registros.', RED_ACCENT),
    ('Supervisor de Almac\u00e9n', 'Aprobaci\u00f3n de usuarios, eliminaci\u00f3n de movimientos y supervisi\u00f3n operativa.', ORANGE),
    ('Supervisor de Operaciones', 'Funciones de supervisi\u00f3n con alcance en operaciones de piso.', ORANGE),
    ('Coordinador de Operaciones', 'Coordinaci\u00f3n general entre \u00e1reas operativas del almac\u00e9n.', ORANGE),
    ('Almacenero', 'Registro de movimientos y consulta de stock. Sin funciones administrativas.', ACCENT),
    ('Operario', 'Registro b\u00e1sico de movimientos e ingreso de datos de manera controlada.', ACCENT),
    ('Auxiliar', 'Acceso de solo lectura y funciones limitadas de registro.', MED_GRAY),
]

y = Inches(1.3)
for title, desc, color in roles:
    add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(0.7), LIGHT_GRAY)
    add_shape_bg(slide, Inches(0.8), y, Inches(0.06), Inches(0.7), color)
    add_textbox(slide, Inches(1.1), y + Inches(0.05), Inches(3), Inches(0.3), title, 14, BLACK, bold=True)
    add_textbox(slide, Inches(1.1), y + Inches(0.35), Inches(10.5), Inches(0.3), desc, 12, MED_GRAY)
    y += Inches(0.8)

add_footer(slide, 5, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 6 — SECCIÓN 03: Movimientos
# ══════════════════════════════════════════════════════════════
make_section_slide('Kardex Racks \u2014 Movimientos',
                   'Registro de ingresos, salidas, devoluciones e insumos no conformes.',
                   GREEN, '03')

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, GREEN)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Pesta\u00f1a Movimientos \u2014 Vista General', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), GREEN)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.7),
            'La pesta\u00f1a \u00abMovimientos\u00bb es el m\u00f3dulo principal para el registro de actividades en el almac\u00e9n. '
            'Dispone de cuatro sub-pesta\u00f1as especializadas seg\u00fan el tipo de operaci\u00f3n. Debajo del formulario '
            'se presenta un registro hist\u00f3rico filtrable de todos los movimientos realizados.',
            14, DARK_GRAY, line_spacing=1.5)

# 4 sub-tabs
tabs = [
    ('Ingreso', 'Registra la entrada de mercader\u00eda a una posici\u00f3n espec\u00edfica del rack. Requiere ubicaci\u00f3n, c\u00f3digo de producto, cantidad, fecha de vencimiento y turno.', GREEN),
    ('Salida', 'Permite dar salida a productos existentes. El sistema busca ubicaciones con stock del c\u00f3digo ingresado y permite salida parcial o total, individual o masiva.', RED_ACCENT),
    ('Devoluci\u00f3n', 'Registra el retorno de mercader\u00eda previamente retirada. Funciona con el mismo formulario que el ingreso, pero genera un tipo de movimiento diferenciado.', ORANGE),
    ('INC', 'Registra insumos no conformes (defectuosos o fuera de especificaci\u00f3n). Incluye un campo adicional para el c\u00f3digo INC, que permite identificar y segregar estos materiales.', ROSE),
]

x = Inches(0.8)
y = Inches(2.3)
for i, (title, desc, color) in enumerate(tabs):
    col = i % 2
    row = i // 2
    cx = x + col * Inches(6.2)
    cy = y + row * Inches(2.3)
    card_bg = add_shape_bg(slide, cx, cy, Inches(5.8), Inches(2.0), LIGHT_GRAY)
    add_shape_bg(slide, cx, cy, Inches(5.8), Inches(0.05), color)
    add_circle_icon(slide, cx + Inches(0.2), cy + Inches(0.2), Inches(0.4), color, str(i+1), 16, WHITE)
    add_textbox(slide, cx + Inches(0.8), cy + Inches(0.2), Inches(4.5), Inches(0.3),
                title, 16, BLACK, bold=True)
    add_textbox(slide, cx + Inches(0.2), cy + Inches(0.75), Inches(5.3), Inches(1.1),
                desc, 12, MED_GRAY, line_spacing=1.4)

add_footer(slide, 6, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 7 — Paso a paso: Ingreso
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, GREEN)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Paso a Paso: Registro de Ingreso', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), GREEN)

steps_ingreso = [
    ('Seleccionar ubicaci\u00f3n', 'Elija el bloque, torre, piso y posici\u00f3n donde se almacenar\u00e1 el producto. Los campos se habilitan en cascada seg\u00fan la selecci\u00f3n anterior.'),
    ('Buscar producto', 'Utilice el buscador del cat\u00e1logo para encontrar el c\u00f3digo del producto. Al seleccionarlo, se autocompletan la descripci\u00f3n y la unidad de medida.'),
    ('Ingresar cantidad', 'Escriba la cantidad a ingresar (valor num\u00e9rico mayor a cero). Puede especificar la fecha de vencimiento o marcar la casilla \u00absin vencimiento\u00bb.'),
    ('Proveedor (si aplica)', 'Para ciertos productos (pel\u00edculas y empaques), el sistema solicita seleccionar un proveedor de la lista disponible.'),
    ('Confirmar registro', 'Si la posici\u00f3n ya contiene stock de otro producto, el sistema muestra una alerta con opci\u00f3n de dar salida al producto existente o confirmar el ingreso.'),
]
y = Inches(1.3)
for i, (title, desc) in enumerate(steps_ingreso):
    add_numbered_step(slide, Inches(0.8), y, i+1, title, desc, GREEN)
    y += Inches(1.15)

add_footer(slide, 7, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 8 — Paso a paso: Salida
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, RED_ACCENT)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Paso a Paso: Registro de Salida', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), RED_ACCENT)

steps_salida = [
    ('Buscar c\u00f3digo de producto', 'Ingrese el c\u00f3digo del producto que desea retirar. El sistema muestra autom\u00e1ticamente todas las ubicaciones donde existe stock de dicho art\u00edculo.'),
    ('Seleccionar ubicaciones', 'Se presenta una lista con las ubicaciones disponibles, ordenadas por FEFO (primer vencido, primera salida). Puede seleccionar una o varias ubicaciones.'),
    ('Definir cantidad a retirar', 'Para cada ubicaci\u00f3n, especifique la cantidad a retirar. Puede usar los botones \u00abParcial\u00bb o \u00abTodo\u00bb para facilitar la operaci\u00f3n.'),
    ('Salida masiva (opcional)', 'Active la selecci\u00f3n m\u00faltiple con el bot\u00f3n \u00abSeleccionar todas\u00bb para procesar la salida en varias ubicaciones de forma simult\u00e1nea.'),
    ('Confirmar operaci\u00f3n', 'Revise el resumen de cantidades y confirme. El sistema valida que exista stock suficiente antes de ejecutar la operaci\u00f3n.'),
]
y = Inches(1.3)
for i, (title, desc) in enumerate(steps_salida):
    add_numbered_step(slide, Inches(0.8), y, i+1, title, desc, RED_ACCENT)
    y += Inches(1.15)

add_footer(slide, 8, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 9 — INC
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ROSE)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Ingreso de Insumo No Conforme (INC)', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), ROSE)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.7),
            'Los insumos no conformes (INC) son aquellos materiales que no cumplen con las especificaciones '
            'de calidad establecidas. RACKLY permite registrarlos desde dos puntos de entrada: la pesta\u00f1a '
            '\u00abMovimientos\u00bb (sub-pesta\u00f1a INC) y la pesta\u00f1a \u00abOcupaci\u00f3n\u00bb (bot\u00f3n INC en cada posici\u00f3n). '
            'Ambos m\u00e9todos almacenan el registro con el campo \u00abcodigo_inc\u00bb para segregarlo del stock regular.',
            14, DARK_GRAY, line_spacing=1.5)

steps_inc = [
    ('Seleccionar ubicaci\u00f3n de destino', 'Determine la posici\u00f3n del rack donde se almacenar\u00e1 el INC. Es recomendable ubicar los materiales no conformes en \u00e1reas segregadas.'),
    ('Buscar producto en el cat\u00e1logo', 'Ingrese el c\u00f3digo del producto. Se autocompletan la descripci\u00f3n y la unidad de medida del cat\u00e1logo maestro.'),
    ('Especificar c\u00f3digo INC', 'Ingrese el identificador \u00fanico del INC (por ejemplo: INC026-120). Este campo es obligatorio y permite rastrear el material defectuoso.'),
    ('Ingresar cantidad y vencimiento', 'Indique la cantidad del insumo no conforme y, si aplica, su fecha de vencimiento. Estos datos quedan registrados en el movimiento.'),
    ('Confirmar registro', 'Al presionar \u00abRegistrar INC\u00bb, el sistema guarda el movimiento con tipo \u00abingreso\u00bb y el c\u00f3digo INC asociado. El material queda visible con un indicador rojo.'),
]
y = Inches(2.3)
for i, (title, desc) in enumerate(steps_inc):
    add_numbered_step(slide, Inches(0.8), y, i+1, title, desc, ROSE)
    y += Inches(1.0)

add_footer(slide, 9, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 10 — SECCIÓN 04: Traslado
# ══════════════════════════════════════════════════════════════
make_section_slide('Kardex Racks \u2014 Traslado',
                   'Transferencia de stock entre ubicaciones del almac\u00e9n.',
                   PURPLE, '04')

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, PURPLE)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Paso a Paso: Traslado de Stock', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), PURPLE)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
            'El m\u00f3dulo de traslado permite mover productos de una ubicaci\u00f3n a otra dentro del almac\u00e9n. '
            'El sistema genera autom\u00e1ticamente una salida en el origen y un ingreso en el destino, manteniendo '
            'la trazabilidad completa de la operaci\u00f3n.',
            14, DARK_GRAY, line_spacing=1.5)

steps_traslado = [
    ('Buscar c\u00f3digo de producto', 'Ingrese el c\u00f3digo del art\u00edculo que desea trasladar. El sistema mostrar\u00e1 las ubicaciones con stock disponible.'),
    ('Seleccionar ubicaci\u00f3n de origen', 'Elija la posici\u00f3n exacta de donde se retirar\u00e1 el producto (bloque, torre, piso y posici\u00f3n).'),
    ('Definir ubicaci\u00f3n de destino', 'Seleccione la nueva posici\u00f3n donde se almacenar\u00e1 el producto trasladado.'),
    ('Especificar cantidad', 'Indique la cantidad de unidades a trasladar. El sistema verificar\u00e1 que exista stock suficiente en el origen.'),
    ('Confirmar traslado', 'El sistema ejecuta la operaci\u00f3n y registra el movimiento en el historial. El stock se actualiza en ambas ubicaciones.'),
]
y = Inches(2.0)
for i, (title, desc) in enumerate(steps_traslado):
    add_numbered_step(slide, Inches(0.8), y, i+1, title, desc, PURPLE)
    y += Inches(1.05)

add_footer(slide, 10, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 11 — Catálogo
# ══════════════════════════════════════════════════════════════
make_section_slide('Kardex Racks \u2014 Cat\u00e1logo',
                   'Gesti\u00f3n del cat\u00e1logo maestro de productos.',
                   ACCENT, '05')

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Gesti\u00f3n del Cat\u00e1logo de Productos', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), ACCENT)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
            'El cat\u00e1logo es el repositorio maestro de todos los c\u00f3digos de productos del almac\u00e9n. '
            'Almacena el c\u00f3digo, la descripci\u00f3n y la unidad de medida de cada art\u00edculo. '
            'Este cat\u00e1logo se utiliza para autocompletar formularios en toda la aplicaci\u00f3n.',
            14, DARK_GRAY, line_spacing=1.5)

steps_catalogo = [
    ('Importar desde Excel', 'Cargue un archivo .xlsx con los datos de los productos. El sistema lee las columnas de c\u00f3digo, descripci\u00f3n y unidad, y las agrega al cat\u00e1logo.'),
    ('Agregar producto manual', 'Ingrese directamente el c\u00f3digo, la descripci\u00f3n y la unidad de medida para crear un nuevo registro en el cat\u00e1logo.'),
    ('B\u00fasqueda y filtro', 'Utilice la barra de b\u00fasqueda para localizar productos existentes por c\u00f3digo o por descripci\u00f3n.'),
    ('Eliminaci\u00f3n (solo admin)', 'Los administradores pueden eliminar registros del cat\u00e1logo. Esta acci\u00f3n no afecta los movimientos ya registrados.'),
]
y = Inches(2.2)
for i, (title, desc) in enumerate(steps_catalogo):
    add_numbered_step(slide, Inches(0.8), y, i+1, title, desc, ACCENT)
    y += Inches(1.15)

add_footer(slide, 11, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 12 — Stock
# ══════════════════════════════════════════════════════════════
make_section_slide('Kardex Racks \u2014 Stock',
                   'Consulta de stock por c\u00f3digo de producto con filtros avanzados.',
                   GREEN, '06')

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, GREEN)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Consulta de Stock', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), GREEN)

steps_stock = [
    ('Buscar c\u00f3digo de producto', 'Ingrese el c\u00f3digo del art\u00edculo para consultar su stock actual. El sistema calcula la diferencia entre ingresos y salidas por ubicaci\u00f3n.'),
    ('Verificar datos del cat\u00e1logo', 'Se muestran los datos del producto: descripci\u00f3n, unidad de medida y stock registrado en el sistema externo Big Magic.'),
    ('Revisar stock por ubicaci\u00f3n', 'El sistema presenta una tabla con el desglose de stock por bloque, torre, piso y posici\u00f3n. Se incluye la fecha de vencimiento y el proveedor.'),
    ('Filtrar por tipo de stock', 'Utilice los botones de filtro: \u00abTodos\u00bb (muestra todo), \u00abDisponibles\u00bb (solo stock normal) y \u00abSolo INC\u00bb (exclusivamente insumos no conformes).'),
]
y = Inches(1.3)
for i, (title, desc) in enumerate(steps_stock):
    add_numbered_step(slide, Inches(0.8), y, i+1, title, desc, GREEN)
    y += Inches(1.3)

add_footer(slide, 12, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 13 — Ocupación
# ══════════════════════════════════════════════════════════════
make_section_slide('Kardex Racks \u2014 Ocupaci\u00f3n',
                   'Mapa visual tridimensional de la ocupaci\u00f3n del almac\u00e9n.',
                   ACCENT, '07')

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Mapa Visual de Ocupaci\u00f3n', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), ACCENT)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
            'La pesta\u00f1a \u00abOcupaci\u00f3n\u00bb muestra una cuadr\u00edcula visual de las posiciones del almac\u00e9n, '
            'organizada por bloques, torres, pisos y posiciones. Cada celda tiene un c\u00f3digo de color que indica '
            'su estado de ocupaci\u00f3n, lo que permite identificar de forma inmediata las \u00e1reas disponibles y las ocupadas.',
            14, DARK_GRAY, line_spacing=1.5)

# Color legend
colors_legend = [
    ('Verde', 'Posici\u00f3n vac\u00eda, disponible para nuevos ingresos.', RGBColor(0x22, 0xC5, 0x5E)),
    ('Azul', 'Posici\u00f3n ocupada por un solo producto.', ACCENT),
    ('Naranja', 'Posici\u00f3n ocupada por m\u00faltiples productos.', ORANGE),
    ('Rosa', 'Posici\u00f3n que contiene insumos no conformes (INC).', ROSE),
]
y = Inches(2.3)
for label, desc, color in colors_legend:
    add_shape_bg(slide, Inches(0.8), y, Inches(0.5), Inches(0.35), color)
    add_textbox(slide, Inches(1.5), y, Inches(1.5), Inches(0.35), label, 13, BLACK, bold=True)
    add_textbox(slide, Inches(3.2), y, Inches(8), Inches(0.35), desc, 12, MED_GRAY)
    y += Inches(0.5)

add_textbox(slide, Inches(0.8), y + Inches(0.3), Inches(11.5), Inches(0.6),
            'Al hacer clic en una posici\u00f3n ocupada, se despliega un panel lateral con el detalle del stock: '
            'c\u00f3digo de producto, descripci\u00f3n, cantidad, fecha de vencimiento y proveedor. Desde ah\u00ed se pueden '
            'ejecutar acciones de ingreso, salida, devoluci\u00f3n, traslado o registro de INC directamente.',
            14, DARK_GRAY, line_spacing=1.5)

add_footer(slide, 13, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 14 — FEFO
# ══════════════════════════════════════════════════════════════
make_section_slide('Kardex Racks \u2014 FEFO',
                   'Control de fechas de vencimiento: primer vencido, primera salida.',
                   ORANGE, '08')

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ORANGE)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Sistema FEFO (First Expired, First Out)', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), ORANGE)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
            'El m\u00f3dulo FEFO implementa la metodolog\u00eda \u00abprimer vencido, primera salida\u00bb, esencial '
            'para la gesti\u00f3n de productos perecederos. El sistema clasifica autom\u00e1ticamente los productos '
            'seg\u00fan su fecha de vencimiento y los presenta con un c\u00f3digo de color intuitivo.',
            14, DARK_GRAY, line_spacing=1.5)

fefo_levels = [
    ('Verde', 'M\u00e1s de 60 d\u00edas para vencer. Producto en \u00f3ptimas condiciones.', RGBColor(0x22, 0xC5, 0x5E), '> 60 d\u00edas'),
    ('Azul', 'Entre 30 y 60 d\u00edas para vencer. Estado de alerta temprana.', ACCENT, '30-60 d\u00edas'),
    ('Naranja', 'Entre 15 y 30 d\u00edas para vencer. Se requiere atenci\u00f3n prioritaria.', ORANGE, '15-30 d\u00edas'),
    ('Rojo', 'Vencido o con menos de 15 d\u00edas. Producto requiere acci\u00f3n inmediata.', RED_ACCENT, '< 15 d\u00edas'),
]

y = Inches(2.2)
for label, desc, color, days in fefo_levels:
    add_shape_bg(slide, Inches(0.8), y, Inches(0.06), Inches(0.65), color)
    add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(0.65), LIGHT_GRAY)
    add_textbox(slide, Inches(1.1), y + Inches(0.05), Inches(1.5), Inches(0.25), label, 14, BLACK, bold=True)
    add_textbox(slide, Inches(1.1), y + Inches(0.33), Inches(8), Inches(0.25), desc, 12, MED_GRAY)
    add_textbox(slide, Inches(10.5), y + Inches(0.15), Inches(1.5), Inches(0.3), days, 13, color, bold=True, alignment=PP_ALIGN.RIGHT)
    y += Inches(0.78)

add_textbox(slide, Inches(0.8), y + Inches(0.3), Inches(11.5), Inches(0.5),
            'Nota: Los insumos no conformes (INC) quedan excluidos del c\u00e1lculo FEFO, ya que se gestionan '
            'de manera independiente seg\u00fan los protocolos de calidad establecidos por la organizaci\u00f3n.',
            13, MED_GRAY, italic=True, line_spacing=1.4)

add_footer(slide, 14, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 15 — Descarga
# ══════════════════════════════════════════════════════════════
make_section_slide('Kardex Racks \u2014 Descarga',
                   'Exportaci\u00f3n de datos a formato Excel (.xlsx).',
                   ACCENT_DARK, '09')

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_DARK)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Exportaci\u00f3n de Datos', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), ACCENT_DARK)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.7),
            'La pesta\u00f1a \u00abDescarga\u00bb permite exportar los datos del sistema a archivos Excel para su '
            'an\u00e1lisis fuera de la plataforma. Se ofrecen dos tipos de exportaci\u00f3n que se ajustan a '
            'diferentes necesidades operativas y administrativas.',
            14, DARK_GRAY, line_spacing=1.5)

# Two export cards
for i, (title, desc, icon_num, color) in enumerate([
    ('Movimientos', 'Genera un archivo Excel con el registro completo de movimientos: tipo, fecha, usuario, c\u00f3digo, '
     'descripci\u00f3n, cantidad, ubicaci\u00f3n, turno, proveedor y, si aplica, el c\u00f3digo INC. Ideal para auditor\u00edas '
     'y reportes operativos.', '1', ACCENT),
    ('Stock', 'Genera un archivo Excel con el estado actual del stock: c\u00f3digo de producto, descripci\u00f3n, '
     'unidad de medida, cantidad total y desglose por ubicaci\u00f3n. Incluye el stock registrado en el sistema '
     'externo Big Magic para comparaci\u00f3n.', '2', GREEN),
]):
    cy = Inches(2.4) + i * Inches(2.2)
    add_shape_bg(slide, Inches(0.8), cy, Inches(11.5), Inches(1.9), LIGHT_GRAY)
    add_shape_bg(slide, Inches(0.8), cy, Inches(0.06), Inches(1.9), color)
    add_circle_icon(slide, Inches(1.2), cy + Inches(0.3), Inches(0.5), color, icon_num, 18, WHITE)
    add_textbox(slide, Inches(2.0), cy + Inches(0.25), Inches(4), Inches(0.35), title, 18, BLACK, bold=True)
    add_textbox(slide, Inches(1.2), cy + Inches(0.9), Inches(10.8), Inches(0.85), desc, 13, MED_GRAY, line_spacing=1.4)

add_footer(slide, 15, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 16 — Usuarios
# ══════════════════════════════════════════════════════════════
make_section_slide('Kardex Racks \u2014 Usuarios',
                   'Administraci\u00f3n de cuentas, roles y permisos del personal.',
                   TEAL, '10')

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, TEAL)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Administraci\u00f3n de Usuarios', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), TEAL)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
            'La pesta\u00f1a \u00abUsuarios\u00bb est\u00e1 disponible exclusivamente para administradores y supervisores. '
            'Permite gestionar el acceso del personal al sistema, asignar roles y mantener el control '
            'de qui\u00e9nes tienen autorizaci\u00f3n para operar.',
            14, DARK_GRAY, line_spacing=1.5)

steps_usuarios = [
    ('Ver solicitudes pendientes', 'Al iniciar, se muestran las cuentas nuevas que a\u00fan no han sido aprobadas. Cada solicitud incluye el nombre y correo del solicitante.'),
    ('Aprobar o denegar acceso', 'Revise los datos del solicitante y apruebe o deniegue su acceso al sistema. Las cuentas denegadas no pueden iniciar sesi\u00f3n.'),
    ('Asignar rol', 'Para cada usuario aprobado, seleccione el rol apropiado de la lista jer\u00e1rquica disponible. El rol define las funciones que el usuario puede realizar.'),
    ('Eliminar usuario', 'En caso necesario, el administrador puede eliminar cuentas del sistema. Esta acci\u00f3n elimina el perfil pero conserva el historial de movimientos registrados.'),
]
y = Inches(2.2)
for i, (title, desc) in enumerate(steps_usuarios):
    add_numbered_step(slide, Inches(0.8), y, i+1, title, desc, TEAL)
    y += Inches(1.15)

add_footer(slide, 16, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 17 — Kardex Piso
# ══════════════════════════════════════════════════════════════
make_section_slide('Kardex Piso',
                   'Operaciones de almac\u00e9n en el \u00e1rea de piso.',
                   RGBColor(0x6D, 0x28, 0xD9), '11')

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, RGBColor(0x6D, 0x28, 0xD9))
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'M\u00f3dulo de Piso de Almac\u00e9n', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), RGBColor(0x6D, 0x28, 0xD9))

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
            'El m\u00f3dulo \u00abKardex Piso\u00bb gestiona las operaciones en el \u00e1rea de piso del almac\u00e9n, '
            'donde se manejan los sectores de almacenamiento a nivel del suelo. Funciona de manera independiente '
            'al Kardex de Racks y cuenta con su propio conjunto de funcionalidades.',
            14, DARK_GRAY, line_spacing=1.5)

piso_tabs = [
    ('Movimientos', 'Historial completo de movimientos realizados en piso. Incluye badges visuales para identificar los insumos no conformes (INC) registrados.', RGBColor(0x6D, 0x28, 0xD9)),
    ('Sectores', 'Vista de los sectores del almac\u00e9n en piso. Permite realizar ingresos, salidas y registros de INC directamente desde cada sector. Cada sector muestra su stock actual.', PURPLE),
    ('Stock', 'Consulta de stock del \u00e1rea de piso por c\u00f3digo de producto. Integra informaci\u00f3n del sistema Big Magic y presenta los datos con control FEFO.', TEAL),
    ('Configuraci\u00f3n', 'Panel administrativo para definir la estructura de sectores y almacenes del \u00e1rea de piso. Solo accesible para personal autorizado.', MED_GRAY),
]

x = Inches(0.8)
y = Inches(2.2)
for i, (title, desc, color) in enumerate(piso_tabs):
    col = i % 2
    row = i // 2
    cx = x + col * Inches(6.2)
    cy = y + row * Inches(2.2)
    add_shape_bg(slide, cx, cy, Inches(5.8), Inches(1.9), LIGHT_GRAY)
    add_shape_bg(slide, cx, cy, Inches(5.8), Inches(0.05), color)
    add_textbox(slide, cx + Inches(0.2), cy + Inches(0.2), Inches(5.3), Inches(0.3),
                title, 16, BLACK, bold=True)
    add_textbox(slide, cx + Inches(0.2), cy + Inches(0.6), Inches(5.3), Inches(1.1),
                desc, 12, MED_GRAY, line_spacing=1.4)

add_footer(slide, 17, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 18 — Funciones Transversales
# ══════════════════════════════════════════════════════════════
make_section_slide('Funciones Transversales',
                   'Modo offline, sincronizaci\u00f3n autom\u00e1tica y actualizaci\u00f3n en tiempo real.',
                   MED_GRAY, '12')

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, MED_GRAY)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Funciones del Sistema', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), MED_GRAY)

features = [
    ('Modo Offline', 'Cuando la conexi\u00f3n a internet se pierde, la aplicaci\u00f3n contin\u00faa funcionando con normalidad. '
     'Los movimientos se almacenan localmente en el navegador (IndexedDB) y se sincronizan autom\u00e1ticamente '
     'al recuperar la conexi\u00f3n. Se genera un identificador \u00fanico (UUID) para garantizar la idempotencia: '
     'si un movimiento se env\u00eda dos veces, el sistema lo detecta y evita duplicados.', PURPLE),
    ('Sincronizaci\u00f3n Autom\u00e1tica', 'El motor de sincronizaci\u00f3n verifica la conectividad con el servidor cada 30 segundos mediante '
     'un ping a la API de Supabase. Al detectar conexi\u00f3n, procesa la cola de movimientos pendientes, env\u00eda '
     'cada uno al servidor y actualiza el cach\u00e9 local con los datos m\u00e1s recientes.', ACCENT),
    ('Actualizaci\u00f3n en Tiempo Real', 'Los cambios realizados por otros usuarios se reflejan de forma instant\u00e1nea gracias a los canales '
     'WebSocket de Supabase Realtime. Si el WebSocket no est\u00e1 disponible, el sistema utiliza un mecanismo '
     'de consulta peri\u00f3dica (polling) cada 8 segundos como respaldo.', GREEN),
    ('Indicador de Conexi\u00f3n', 'La barra superior de la aplicaci\u00f3n muestra un indicador visual del estado de conexi\u00f3n: '
     'verde (en l\u00ednea), amarillo (sincronizando) o rojo (sin conexi\u00f3n). Tambi\u00e9n muestra el conteo '
     'de movimientos pendientes por sincronizar.', ORANGE),
]

y = Inches(1.3)
for title, desc, color in features:
    add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(1.25), LIGHT_GRAY)
    add_shape_bg(slide, Inches(0.8), y, Inches(0.06), Inches(1.25), color)
    add_textbox(slide, Inches(1.1), y + Inches(0.08), Inches(10.8), Inches(0.3), title, 15, BLACK, bold=True)
    add_textbox(slide, Inches(1.1), y + Inches(0.42), Inches(10.8), Inches(0.75), desc, 12, MED_GRAY, line_spacing=1.4)
    y += Inches(1.4)

add_footer(slide, 18, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 19 — Referencias APA
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, MED_GRAY)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
            'Referencias', 28, BLACK, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.2), MED_GRAY)

references = [
    'Hidalgo, M. (2026). RACKLY: Sistema de gesti\u00f3n de almacenes (Versi\u00f3n 2.0) [Software inform\u00e1tico]. https://rackly.pages.dev',
    'Real Academia Espa\u00f1ola. (2024). Diccionario de la lengua espa\u00f1ola (24.б╙ ed.). https://dle.rae.es',
    'American Psychological Association. (2020). Publication manual of the American Psychological Association (7.б╙ ed.). https://doi.org/10.1037/0000165-000',
    'Supabase. (2026). Supabase: The open source Firebase alternative. https://supabase.com',
    'Vercel. (2026). Next.js: The React framework for the web. https://nextjs.org',
    'Monge, F. y Arranz, P. (2023). Gesti\u00f3n de almacenes: Manual pr\u00e1ctico de log\u00edstica industrial. Ediciones Paraninfo.',
    'Real Academia Espa\u00f1ola y Asociaci\u00f3n de Academias de la Lengua Espa\u00f1ola. (2010). Ortograf\u00eda de la lengua espa\u00f1ola. Espasa.',
]

y = Inches(1.3)
for ref in references:
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.5), Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ref
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_MAIN
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = Pt(16)
    # Hanging indent effect
    p.first_line_indent = Inches(-0.4)
    p.left_indent = Inches(0.4)
    y += Inches(0.72)

add_footer(slide, 19, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════
# DIAPOSITIVA 20 — Cierre
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
add_shape_bg(slide, Inches(0), SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), ACCENT)

add_circle_icon(slide, Inches(5.9), Inches(1.8), Inches(1.5), ACCENT, 'R', 52, WHITE)

add_textbox(slide, Inches(2), Inches(3.6), Inches(9.3), Inches(0.7),
            'RACKLY', 44, BLACK, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(0.4),
            'Sistema de Gesti\u00f3n de Almacenes v2.0', 18, MED_GRAY, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(5.5), Inches(4.8), Inches(2.3), LIGHT_GRAY)
add_textbox(slide, Inches(2), Inches(5.2), Inches(9.3), Inches(0.4),
            'Manual de Usuario \u2014 Junio 2026', 14, MED_GRAY, italic=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(5.7), Inches(9.3), Inches(0.4),
            'Desarrollado por Miguel Hidalgo', 13, MED_GRAY, alignment=PP_ALIGN.CENTER)

# ─── Guardar ───
output_path = '/home/z/my-project/download/RACKLY_Manual_de_Usuario.pptx'
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f'Presentaci\u00f3n guardada en: {output_path}')
print(f'Total de diapositivas: {len(prs.slides)}')
